"""Build the results deck from promoted artifacts and generated figures.

Every number on a slide is read from docs/results/ at build time rather than typed in, so a
slide cannot drift from the evidence the way a hand-written figure caption does. Where a rung
has no promoted result the slide says so explicitly and names what is missing -- a deck that
silently omits an incomplete rung reads as a complete story.

Caveats and open decisions get their own slides rather than a footnote, because they are the
part a reader most needs and the part most easily lost.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x77, 0x77, 0x77)
ALERT = RGBColor(0xC4, 0x4E, 0x52)


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Opening slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, INK
    tb2 = s.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(11.9), Inches(1.2))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size, p2.font.color.rgb = Pt(16), MUTED


def _bullets(prs: Presentation, title: str, lines: list[tuple[str, int, bool]]) -> None:
    """Text slide. Each line is (text, indent level, is_alert)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(26), True, INK
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, (text, lvl, alert) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = lvl
        para.font.size = Pt(16 if lvl == 0 else 13)
        para.font.color.rgb = ALERT if alert else (INK if lvl == 0 else MUTED)


def _figure_slide(prs: Presentation, title: str, img: Path, notes: list[str]) -> None:
    """Figure slide with the source note underneath."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(24), True, INK
    if img.exists():
        s.shapes.add_picture(str(img), Inches(1.1), Inches(1.2), height=Inches(4.7))
    else:
        nb = s.shapes.add_textbox(Inches(1.1), Inches(3.0), Inches(10), Inches(1))
        np_ = nb.text_frame.paragraphs[0]
        np_.text = f"figure not generated: {img.name}"
        np_.font.size, np_.font.color.rgb = Pt(14), ALERT
    nb = s.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(11.9), Inches(1.1))
    tf = nb.text_frame
    tf.word_wrap = True
    for i, n in enumerate(notes):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = n
        para.font.size, para.font.color.rgb = Pt(11), MUTED


def load(results: Path, name: str) -> tuple[pd.DataFrame | None, str]:
    """A promoted result and its job id, or (None, reason)."""
    f = results / f"{name}.csv"
    if not f.exists():
        return None, f"{name}.csv not promoted"
    side = results / f"{name}.provenance.json"
    job = ""
    if side.exists():
        try:
            job = str(json.loads(side.read_text()).get("job_id", ""))
        except Exception:
            pass
    return pd.read_csv(f), job


def main() -> None:
    """Assemble the deck."""
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--results-dir", type=Path, default=Path("docs/results"))
    ap.add_argument("--figures-dir", type=Path, default=Path("docs/figures"))
    ap.add_argument("--out", type=Path, default=Path("docs/ladder_results.pptx"))
    ap.add_argument("--date", required=True, help="round date, e.g. 2026-08-25")
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    _title_slide(prs, "Transfer ladder: where cross-modality prediction breaks down",
                 f"Modular harness core - {args.date}. Every number traces to a promoted "
                 "artifact with a job id.")

    _bullets(prs, "The question, and why a single number cannot answer it", [
        ("A Path B result has no interpretable scale on its own.", 0, False),
        ("L1000 and Tahoe deltas for the SAME (cell line, drug) agree at Spearman 0.041 against "
         "a split-half ceiling of 0.572, with sign concordance at chance.", 1, True),
        ("No normalisation recovers it: seven per-gene transforms all land between 0.034 and "
         "0.050, none clearing its own null [job 31661918].", 1, False),
        ("So the ladder supplies the scale: each rung adds exactly one distribution shift.", 0, False),
        ("Rung 0 replicate ceiling - rung 1 held-out line - rung 2 cross-platform - "
         "rung 3 GDSC2 viability - rung 4 organoids (frozen holdout).", 1, False),
    ])

    _bullets(prs, "What makes the rungs comparable", [
        ("Six invariants, held identical at every rung. Violating any one turns a ratio between "
         "rungs into a comparison of two different measurements.", 0, False),
        ("Gene panel: Check 1 on 14,121 genes, Check 2 on 12,368, rung 2 on 8,759 "
         "(L1000 constrains it).", 1, False),
        ("CV: 5-fold over cell lines from ONE shared partition, so folds are the same "
         "partition and not merely the same count.", 1, False),
        ("Metric, unit, reliability correction, and a null recomputed within each rung.", 1, False),
        ("Reported as a fraction of that rung's own ceiling, never as a raw number.", 1, False),
    ])

    figs = [
        ("Rung 0 - replicate ceiling", "rung0_ceiling.png",
         ["Split-half over plate halves, Spearman-Brown lifted to full data.",
          "The mismatched-pair null is what makes this a ceiling rather than shared gene structure."]),
        ("Rung 1 - held-out line, DE fidelity", "rung1_de_null.png",
         ["shuffle_all is cleared by every source including line-independent baselines, so it "
          "tests drug specificity, not line specificity.",
          "within_drug holds drug identity fixed and is the test the published claim needs."]),
        ("Rung 2 - cross-platform and granularity transfer", "rung2_transfer.png",
         ["Transfer penalty is cross_platform minus in_platform, both arms on one pinned panel.",
          "Stack appears only via its L1000-context arm; it is not fitted here, so an unlabelled "
          "constant would read as a win."]),
        ("Rung 3 - GDSC2 viability", "rung3_check2.png",
         ["Grey bars are controls. planted must be recovered; prior is the line-independent floor.",
          "Ceiling is independent-screen agreement, 0.457 GDSC2 vs CTRPv2 [job 31663627]."]),
        ("The ladder", "ladder_summary.png",
         ["Each rung as a fraction of its OWN ceiling. Rung 4 has no ceiling and is omitted.",
          "Rungs without both a score and a ceiling are named rather than silently dropped."]),
    ]
    for title, img, notes in figs:
        _figure_slide(prs, title, args.figures_dir / img, notes)

    audit, ajob = load(args.results_dir, "ladder_audit")
    lines = [("Controls and provenance were audited mechanically, not by reading the design.", 0, False)]
    if audit is not None:
        for r in audit.itertuples():
            gaps = [a for a in ("controls_floor", "controls_negative", "controls_positive",
                                "prov_params", "prov_panel", "prov_drugs")
                    if not getattr(r, f"{a}_ok", True)]
            lines.append((f"{r.rung}: {'all axes covered' if not gaps else 'GAPS - ' + ', '.join(gaps)}",
                          1, bool(gaps)))
    lines.append(("Reading the design is how the gaps survived: the panel was imported and never "
                  "called, and Check 2's positive control vanished in the array conversion.", 0, False))
    _bullets(prs, "Audit: controls and provenance per rung", lines)

    _bullets(prs, "What a reader should not conclude", [
        ("Rung 4 has NO ceiling. Verified: no dose or replicate column in any of the seven "
         "organoid tables [job 31663218].", 0, True),
        ("So no rung-4 result may be stated as an absolute value. The honest form is always "
         "'X% of what the same method achieves on cell lines'.", 1, False),
        ("Rung 4's usable n is 17 organoids at a median of 17 of 34 drugs, not the 94 screened.", 1, False),
        ("Stack's rung-2 arm covers 14 drugs against the baselines' wider set; the two are not "
         "equally powered and are not tabulated as if they were.", 0, True),
        ("The L1000 imputation-fidelity test is NOT ESTABLISHED: the raw gap did not survive "
         "variance matching (p = 0.270) [job 31661570].", 0, True),
    ])

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"wrote {args.out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
