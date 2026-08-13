"""Rebuild the results slides (5-7) of ``docs/prospective_evaluation_harness_overview.pptx``.

Content is the reported values in ``docs/tahoe_generation_results.md``. The script is
idempotent: it strips the body shapes it owns on those slides and re-adds them, so it can be
re-run after the results doc changes. Slide 6's figure comes from
``scripts/plot_generation_eval_summary.py`` -- run that first.

    python3 scripts/plot_generation_eval_summary.py
    python3 scripts/update_harness_overview_slides.py
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

DECK = Path("docs/prospective_evaluation_harness_overview.pptx")
FIGURE = Path("docs/figures/generation_eval_summary.png")

HEADER_BG = RGBColor(0x2C, 0x3E, 0x50)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
BODY_FG = RGBColor(0x22, 0x22, 0x22)
ROW_A, ROW_B = RGBColor(0xF2, 0xF5, 0xF8), RGBColor(0xFF, 0xFF, 0xFF)
RED_BG, GREEN_BG = RGBColor(0xF6, 0xDD, 0xDB), RGBColor(0xD8, 0xEE, 0xDD)
RED_FG, GREEN_FG = RGBColor(0x9B, 0x2C, 0x22), RGBColor(0x1E, 0x6B, 0x3A)
CAPTION_FG = RGBColor(0x44, 0x44, 0x44)
TABLE_STYLE = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
CELL_MARGIN = Emu(54864)

# --- Check 1 ------------------------------------------------------------------------------
CHECK1_HEAD = ["Source", "r  (predicted vs real)", "r  off-diagonal", "specificity rank", "pairs"]
CHECK1_ROWS = [
    ("additive", "0.225", "0.095", "0.885", "1600", None),
    ("nmf", "0.221", "0.088", "0.912", "1600", None),
    ("pca", "0.207", "0.083", "0.896", "1600", None),
    ("knn", "0.178", "0.067", "0.904", "1600", None),
    ("stack (gen, cytokine-aligned)", "0.012", "-0.002", "0.644", "1568", "red"),
    ("stack (gen, drug-aligned, unfiltered)", "0.021", "0.006", "0.665", "1568", "red"),
    ("stack (gen, drug-aligned, leak-excluded)", "0.021", "0.006", "0.665", "1563", "red"),
]
CHECK1_TAKEAWAY = (
    "Stack generation is null: r = 0.012, off-diagonal ~ 0, specificity rank 0.64 (random ~ 0.5) — "
    "below the line-independent additive floor (0.225) and far below the 0.46 reproducibility ceiling."
)

# --- Hallmark gate ------------------------------------------------------------------------
GATE_HEAD = ["Hallmark set", "interaction", "global", "p vs random", "clears gate?"]
GATE_ROWS = [
    ("P53 pathway", "0.009", "-0.091", "0.645", "fail", "red"),
    ("Apoptosis", "0.004", "0.003", "0.720", "fail", "red"),
    ("E2F targets", "0.032", "0.049", "0.025", "pass", "green"),
    ("G2M checkpoint", "0.047", "0.117", "0.000", "pass", "green"),
]
GATE_TAKEAWAY = (
    "Only the two proliferation sets beat random; the cell-death sets are indistinguishable from "
    "random on Tahoe, so a death-signature readout is underpowered here."
)

# --- Check 2, fixed-signature readouts ----------------------------------------------------
SIG_HEAD = ["source", "method", "global", "interaction", "per-drug", "sel. gap@1", "sel. gap@3", "p_label"]
SIG_ROWS = [
    ("additive", "hallmark", "0.088", "-0.063", "-0.033", "0.855", "0.579", "0.979", None),
    ("additive", "proliferation", "0.097", "-0.016", "-0.035", "0.585", "0.547", "0.884", None),
    ("knn", "hallmark", "0.080", "0.022", "-0.008", "0.788", "0.587", "0.220", None),
    ("knn", "proliferation", "0.092", "0.015", "0.063", "0.628", "0.415", "0.360", None),
    ("pca", "hallmark", "0.089", "-0.075", "0.053", "0.855", "0.579", "0.985", None),
    ("pca", "proliferation", "0.123", "-0.010", "0.049", "0.585", "0.535", "0.848", None),
    ("nmf", "hallmark", "0.077", "-0.068", "-0.072", "0.855", "0.579", "0.981", None),
    ("nmf", "proliferation", "0.094", "-0.016", "-0.023", "0.585", "0.549", "0.882", None),
    ("stack (gen, cytokine-aligned)", "hallmark", "-0.128", "-0.030", "-0.090", "0.733", "0.581", "0.809", "red"),
    ("stack (gen, cytokine-aligned)", "proliferation", "-0.150", "-0.014", "-0.090", "0.849", "0.674", "0.560", "red"),
    ("stack (gen, drug-aligned, unfiltered)", "hallmark", "-0.077", "0.019", "-0.084", "0.687", "0.532", "0.322", "red"),
    ("stack (gen, drug-aligned, unfiltered)", "proliferation", "-0.038", "0.081", "-0.057", "0.317", "0.209", "0.009", "red"),
    ("stack (gen, drug-aligned, leak-excluded)", "hallmark", "-0.074", "0.021", "-0.079", "0.687", "0.532", "0.296", "red"),
    ("stack (gen, drug-aligned, leak-excluded)", "proliferation", "-0.035", "0.082", "-0.057", "0.317", "0.208", "0.004", "red"),
]

# --- Check 2, trained penalized ladder ----------------------------------------------------
LADDER_HEAD = [
    "representation",
    "L2  global / int",
    "L1  global / int",
    "EN  global / int",
    "per-drug (L2)",
    "p_label (L2)",
    "sel. gap@1",
    "sel. gap@3",
]
LADDER_ROWS = [
    ("expr", "0.475 / -0.037", "0.599 / -0.108", "0.604 / -0.123", "-0.011", "0.863", "0.354", "0.119", None),
    ("additive", "0.628 / -0.095", "0.603 / -0.159", "0.601 / -0.155", "-0.149", "0.997", "0.264", "0.091", None),
    ("knn", "0.547 / -0.068", "0.617 / -0.171", "0.618 / -0.169", "-0.067", "0.975", "0.250", "0.101", None),
    ("pca", "0.585 / +0.007", "0.634 / -0.108", "0.634 / -0.103", "+0.018", "0.498", "0.219", "0.102", None),
    ("nmf", "0.550 / +0.007", "0.610 / -0.198", "0.614 / -0.178", "-0.101", "0.596", "0.251", "0.082", None),
    ("stack (gen, cytokine-aligned)", "0.539 / -0.003", "0.568 / -0.182", "0.574 / -0.164", "-0.071", "0.625", "0.320", "0.144", "red"),
    ("stack (gen, drug-aligned, unfiltered)", "0.561 / -0.082", "0.575 / -0.145", "0.570 / -0.147", "-0.135", "0.987", "0.343", "0.133", "red"),
    ("stack (gen, drug-aligned, leak-excluded)", "0.561 / -0.079", "0.572 / -0.134", "0.562 / -0.123", "-0.140", "0.987", "0.324", "0.122", "red"),
    ("base (embed)", "0.644 / +0.119", "0.612 / -0.166", "0.613 / -0.170", "+0.200", "0.001", "0.273", "0.102", "green"),
    ("aligned (embed)", "0.618 / +0.045", "0.623 / -0.097", "0.625 / -0.103", "+0.059", "0.175", "0.240", "0.096", None),
]
LADDER_TAKEAWAY = (
    "The base (unaligned) Stack embedding is the only representation that captures cell-line-specific "
    "response: interaction +0.119, per-drug +0.200, p_label 0.001 under ridge, where expression, PCA, NMF "
    "and every generated delta sit at ~0 and non-significant. The signal is dense — L1/EN sparsify it away "
    "(interaction ~ -0.17) — and cytokine alignment does not transfer (base > aligned on every interaction "
    "metric). Overall potency (~0.6) is solved by everything; the interaction ceiling is 0.31-0.47."
)


def _clear_body(slide, keep_banner: bool = True) -> None:
    """Drop every shape except the title banner, so a re-run does not stack duplicates."""
    for shape in list(slide.shapes):
        is_banner = keep_banner and shape.name.startswith("Rectangle")
        if not is_banner:
            shape._element.getparent().remove(shape._element)


def _set_banner(slide, text: str) -> None:
    banner = next(s for s in slide.shapes if s.name.startswith("Rectangle"))
    tf = banner.text_frame
    para = tf.paragraphs[0]
    template = para.runs[0]
    for run in list(para.runs)[1:]:
        run._r.getparent().remove(run._r)
    template.text = text
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)


def _textbox(slide, text, *, left, top, width, height, size, bold=False, color=BODY_FG):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _style_cell(cell, text, *, size, bold, fg, bg, align):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = CELL_MARGIN
    cell.margin_top = cell.margin_bottom = Emu(0)
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg


def _add_table(slide, header, rows, *, left, top, widths, row_height, head_size, body_size):
    """Deck-styled table: dark header, zebra body, optional per-row pass/fail tint."""
    shape = slide.shapes.add_table(
        len(rows) + 1, len(header), Inches(left), Inches(top), Inches(sum(widths)), Inches(row_height * (len(rows) + 1))
    )
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    style = table._tbl.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    if style is not None:
        style.text = TABLE_STYLE
    for col, width in zip(table.columns, widths):
        col.width = Inches(width)
    for row in table.rows:
        row.height = Inches(row_height)

    for ci, label in enumerate(header):
        _style_cell(
            table.cell(0, ci),
            label,
            size=head_size,
            bold=True,
            fg=HEADER_FG,
            bg=HEADER_BG,
            align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
        )
    for ri, row in enumerate(rows):
        *values, tint = row
        zebra = ROW_A if ri % 2 == 0 else ROW_B
        bg = {"red": RED_BG, "green": GREEN_BG, None: zebra}[tint]
        for ci, value in enumerate(values):
            fg = BODY_FG
            if tint == "red" and value in {"fail", "pass"}:
                fg = RED_FG
            elif tint == "green" and value in {"fail", "pass"}:
                fg = GREEN_FG
            _style_cell(
                table.cell(ri + 1, ci),
                value,
                size=body_size,
                bold=(ci == 0) or tint is not None,
                fg=fg,
                bg=bg,
                align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
            )
    return shape


def build_slide5(slide) -> None:
    _clear_body(slide)
    _set_banner(slide, "Results tables — generation quality and the Hallmark gate")
    _textbox(
        slide,
        "How well each source reproduces the real Tahoe change "
        "(1,600 line-drug pairs; 1,568 for Stack — top 2,000 variable genes)",
        left=0.40, top=1.02, width=12.50, height=0.31, size=12.5, bold=True,
    )
    _add_table(
        slide, CHECK1_HEAD, CHECK1_ROWS,
        left=0.40, top=1.42, widths=[2.3, 2.2, 1.9, 1.9, 0.9], row_height=0.34,
        head_size=14, body_size=16,
    )
    _textbox(slide, CHECK1_TAKEAWAY, left=0.40, top=3.52, width=12.50, height=0.60, size=12.5, color=CAPTION_FG)
    _textbox(
        slide,
        "Is the readout powered? The real Tahoe change scored through Hallmark sets vs random gene sets",
        left=0.40, top=4.20, width=12.50, height=0.31, size=12.5, bold=True,
    )
    _add_table(
        slide, GATE_HEAD, GATE_ROWS,
        left=0.40, top=4.58, widths=[2.6, 1.65, 1.5, 1.75, 1.7], row_height=0.34,
        head_size=12, body_size=14,
    )
    _textbox(slide, GATE_TAKEAWAY, left=0.40, top=6.42, width=12.50, height=0.60, size=12.5, color=CAPTION_FG)


def build_slide6(slide) -> None:
    if not FIGURE.exists():
        raise SystemExit(f"missing {FIGURE}; run scripts/plot_generation_eval_summary.py first")
    _clear_body(slide)
    _set_banner(
        slide,
        "Current results — generation null; the base Stack embedding is the exception",
    )
    from PIL import Image  # noqa: PLC0415 -- only needed for the figure aspect ratio

    with Image.open(FIGURE) as img:
        aspect = img.width / img.height
    top, max_width, max_height = 1.05, 13.0, 6.30
    width = min(max_width, max_height * aspect)
    height = width / aspect
    slide.shapes.add_picture(
        str(FIGURE), Inches((13.333 - width) / 2), Inches(top), Inches(width), Inches(height)
    )


def build_slide7(slide) -> None:
    _clear_body(slide)
    _set_banner(slide, "Results table — predicting GDSC2 viability, grouped 5-fold by cell line")
    _textbox(
        slide,
        "global = overall drug potency · interaction = cell-line-specific response · "
        "per-drug = within-drug line ranking · sel. gap@k = best-drug shortlisting, lowest across the "
        "penalty sweep (lower = better) · p_label = personalization null (small = real signal)",
        left=0.40, top=0.95, width=12.60, height=0.40, size=10, color=CAPTION_FG,
    )
    _textbox(
        slide,
        "Fixed-signature readouts on the delta sources  (n = 1,313 pairs)",
        left=0.40, top=1.40, width=12.50, height=0.28, size=11, bold=True,
    )
    _add_table(
        slide, SIG_HEAD, SIG_ROWS,
        left=0.40, top=1.70, widths=[1.55, 1.55, 1.35, 1.55, 1.35, 1.40, 1.40, 1.35], row_height=0.20,
        head_size=8.5, body_size=8.5,
    )
    _textbox(
        slide,
        "Trained penalized models — representation-controlled  (n = 1,303 pairs; L2 = ridge, "
        "L1 = lasso, EN = elastic-net; per-drug and p_label from the L2 fit; folds grouped by cell line, so no line is in both train and test)",
        left=0.40, top=4.00, width=12.50, height=0.28, size=11, bold=True,
    )
    _add_table(
        slide, LADDER_HEAD, LADDER_ROWS,
        left=0.40, top=4.30, widths=[1.85, 1.80, 1.80, 1.80, 1.15, 1.15, 0.98, 0.97], row_height=0.22,
        head_size=9, body_size=9,
    )
    _textbox(slide, LADDER_TAKEAWAY, left=0.40, top=6.36, width=12.50, height=0.90, size=11, color=CAPTION_FG)


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--deck", type=Path, default=DECK)
    args = ap.parse_args()

    prs = Presentation(args.deck)
    build_slide5(prs.slides[4])
    build_slide6(prs.slides[5])
    build_slide7(prs.slides[6])
    prs.save(args.deck)
    print(f"updated slides 5-7 of {args.deck}")


if __name__ == "__main__":
    main()
